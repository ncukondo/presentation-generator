#!/usr/bin/env python3
"""
昨年度スライド(.pptx)から「素材」をスライド単位で抽出し、ファイルとして整理する。

抽出するもの:
  1. テキスト・構成   -> materials/<year>/text/slideNN.md  +  INDEX.md
  2. 画像・図(media) -> materials/<year>/media/slideNN_imgMM.ext  +  一覧
  3. マニフェスト     -> materials/<year>/manifest.json (機械可読)

スライド画像(PNG)は本スクリプトでは生成しない(python-pptx はレンダリング不可)。
PNG 化は tools/export_pptx_png.ts (PowerPoint COM) を使う。

Usage:
  python3 tools/extract_materials.py \
      --src "../2025年度/2025年版_医学教育と研究.pptx" \
      --year 2025
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu


def iter_shapes(shapes):
    """グループ図形を再帰的に展開して全 shape を列挙する。"""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    lines = []
    for para in shape.text_frame.paragraphs:
        text = "".join(run.text for run in para.runs)
        if text.strip():
            indent = "  " * (para.level or 0)
            lines.append(f"{indent}{text}")
    return "\n".join(lines)


def table_markdown(shape) -> str:
    tbl = shape.table
    rows = []
    for r in tbl.rows:
        cells = [c.text.replace("\n", " ").strip() for c in r.cells]
        rows.append("| " + " | ".join(cells) + " |")
    if not rows:
        return ""
    header_sep = "| " + " | ".join("---" for _ in tbl.columns) + " |"
    return "\n".join([rows[0], header_sep, *rows[1:]])


def emu_to_cm(v) -> float:
    return round(Emu(v).cm, 1) if v is not None else 0.0


def extract(src: Path, out_base: Path, year: str) -> dict:
    prs = Presentation(str(src))
    media_dir = out_base / "media"
    text_dir = out_base / "text"
    media_dir.mkdir(parents=True, exist_ok=True)
    text_dir.mkdir(parents=True, exist_ok=True)

    manifest = {"source": src.name, "year": year, "slide_count": len(prs.slides), "slides": []}
    index_lines = [
        f"# {year}年版 スライド素材インデックス",
        "",
        f"- ソース: `{src.name}`",
        f"- スライド数: {len(prs.slides)}",
        "",
        "| # | 見出し（推定） | テキスト塊 | 画像 |",
        "|---|---|---|---|",
    ]

    for idx, slide in enumerate(prs.slides, start=1):
        nn = f"{idx:02d}"
        texts: list[str] = []
        tables_md: list[str] = []
        images: list[dict] = []
        title_guess = ""

        img_no = 0
        for shape in iter_shapes(slide.shapes):
            # title 推定
            if shape == slide.shapes.title and shape.has_text_frame:
                title_guess = shape.text_frame.text.strip().replace("\n", " ")

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.shape_type == MSO_SHAPE_TYPE.LINKED_PICTURE:
                try:
                    image = shape.image
                except Exception:
                    continue
                img_no += 1
                ext = image.ext
                fname = f"slide{nn}_img{img_no:02d}.{ext}"
                (media_dir / fname).write_bytes(image.blob)
                images.append({
                    "file": fname,
                    "ext": ext,
                    "w_cm": emu_to_cm(shape.width),
                    "h_cm": emu_to_cm(shape.height),
                    "alt": (shape.name or "").strip(),
                })
                continue

            if getattr(shape, "has_table", False):
                md = table_markdown(shape)
                if md:
                    tables_md.append(md)
                continue

            t = shape_text(shape)
            if t.strip():
                texts.append(t)

        # notes
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        # per-slide markdown
        md_lines = [f"# Slide {idx}", ""]
        if title_guess:
            md_lines += [f"**見出し:** {title_guess}", ""]
        if texts:
            md_lines += ["## テキスト", ""]
            for block in texts:
                md_lines += [block, ""]
        if tables_md:
            md_lines += ["## 表", ""]
            for md in tables_md:
                md_lines += [md, ""]
        if images:
            md_lines += ["## 画像", ""]
            for im in images:
                md_lines.append(f"- `media/{im['file']}` ({im['w_cm']}×{im['h_cm']}cm)")
            md_lines.append("")
        if notes:
            md_lines += ["## ノート", "", notes, ""]
        (text_dir / f"slide{nn}.md").write_text("\n".join(md_lines), encoding="utf-8")

        manifest["slides"].append({
            "index": idx,
            "title": title_guess,
            "text_blocks": texts,
            "tables": tables_md,
            "images": images,
            "notes": notes,
        })

        text_summary = (texts[0][:18] + "…") if texts else ""
        index_lines.append(
            f"| {idx} | {title_guess or '—'} | {len(texts)}塊 {text_summary} | {len(images)} |"
        )

    (out_base / "INDEX.md").write_text("\n".join(index_lines) + "\n", encoding="utf-8")
    (out_base / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return manifest


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--src", required=True, help="昨年度 pptx へのパス")
    ap.add_argument("--year", default="2025")
    ap.add_argument("--out", default=None, help="出力ベース(既定: materials/<year>)")
    args = ap.parse_args()

    src = Path(args.src).expanduser().resolve()
    if not src.exists():
        raise SystemExit(f"ソースが見つかりません: {src}")
    root = Path(__file__).resolve().parent.parent
    out_base = Path(args.out).resolve() if args.out else root / "materials" / args.year

    m = extract(src, out_base, args.year)
    total_imgs = sum(len(s["images"]) for s in m["slides"])
    print(f"Done: {m['slide_count']} slides, {total_imgs} images -> {out_base}")


if __name__ == "__main__":
    main()
